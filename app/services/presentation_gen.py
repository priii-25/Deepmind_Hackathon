"""
Presentation generation service — Gemini 2.5 Flash Image (nanobanana) + python-pptx.

Generates slide images using Gemini, then assembles them into a downloadable PPTX.
"""

import asyncio
import io
import logging
import os
import tempfile
from typing import Optional

logger = logging.getLogger(__name__)

# ── Gemini client singleton (reuses pattern from fashion_photo.py) ────

_gemini_client = None


def _get_gemini_client():
    """Lazy-load and cache the google-genai client."""
    global _gemini_client
    if _gemini_client is not None:
        return _gemini_client
    try:
        from google import genai
        from ..core.config import get_settings

        settings = get_settings()
        api_key = settings.gemini_api_key
        if not api_key:
            raise ValueError("GEMINI_API_KEY is required for presentation generation")
        _gemini_client = genai.Client(api_key=api_key)
        return _gemini_client
    except ImportError:
        raise ImportError(
            "google-genai package is required for presentation generation. "
            "Install it with: pip install google-genai"
        )


# ── Slide prompt builders ─────────────────────────────────────────────

SLIDE_STYLE = (
    "Clean, modern, professional corporate presentation design. "
    "Dark gradient background (deep navy blue to dark blue). "
    "White and light-colored text. Subtle geometric accent shapes. "
    "Widescreen 16:9 aspect ratio."
)


def build_title_slide_prompt(topic: str) -> str:
    """Build a prompt that tells Gemini to generate a TITLE SLIDE image."""
    return f"""Generate an image of a professional PowerPoint presentation title slide.

=== CRITICAL INSTRUCTIONS ===
- You MUST generate an image that looks EXACTLY like a real PowerPoint slide.
- This is the SLIDE ITSELF — NOT a photograph, NOT an illustration about the topic.
- The image should look like a screenshot of a PowerPoint slide.

=== SLIDE DETAILS ===
- Slide type: TITLE SLIDE (Slide 1 of 2)
- Title text: "{topic}"
- Include a subtle tagline or subtitle area below the title
- Design: {SLIDE_STYLE}
- The title should be large, bold, and centered
- Add subtle decorative elements (thin lines, geometric shapes, gradient overlays)
- Include a small footer area at the bottom

=== IMPORTANT ===
Generate the slide itself as it would appear in PowerPoint presentation mode.
The output must be a single widescreen (16:9) slide image with text rendered ON the slide."""


def build_content_slide_prompt(topic: str) -> str:
    """Build a prompt that tells Gemini to generate a CONTENT SLIDE image."""
    return f"""Generate an image of a professional PowerPoint presentation content slide.

=== CRITICAL INSTRUCTIONS ===
- You MUST generate an image that looks EXACTLY like a real PowerPoint slide.
- This is the SLIDE ITSELF — NOT a photograph, NOT an illustration about the topic.
- The image should look like a screenshot of a PowerPoint slide.

=== SLIDE DETAILS ===
- Slide type: CONTENT SLIDE (Slide 2 of 2)
- Topic: "{topic}"
- Include a clear header/title at the top of the slide
- Include 3-4 bullet points with concise, informative text about the topic
- Each bullet point should have a short heading and a brief explanation
- Design: {SLIDE_STYLE}
- Use consistent styling that matches a title slide
- Add small icons or visual markers next to each bullet point
- Include a subtle footer area

=== IMPORTANT ===
Generate the slide itself as it would appear in PowerPoint presentation mode.
The output must be a single widescreen (16:9) slide image with text rendered ON the slide.
The bullet points must contain REAL, RELEVANT content about "{topic}"."""


# ── Sync Gemini call (runs in executor) ───────────────────────────────

def _sync_generate_slide_image(prompt: str) -> Optional[bytes]:
    """
    Call Gemini 2.5 Flash Image (nanobanana) to generate a single slide image.

    Returns PNG bytes or None if generation fails.
    """
    from google.genai import types

    client = _get_gemini_client()

    response = client.models.generate_content(
        model="gemini-2.5-flash-image",
        contents=prompt,
        config=types.GenerateContentConfig(
            response_modalities=["IMAGE"],
        ),
    )

    # Extract image from response
    for part in response.candidates[0].content.parts:
        if part.inline_data is not None:
            # Save to temp file and read back as bytes
            image = part.as_image()
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                image.save(tmp.name)
                tmp_path = tmp.name
            with open(tmp_path, "rb") as f:
                image_bytes = f.read()
            os.remove(tmp_path)
            return image_bytes

    logger.warning("Gemini returned no image for slide generation")
    return None


# ── PPTX assembly ─────────────────────────────────────────────────────

def _sync_assemble_pptx(slide_images: list[bytes], title: str) -> bytes:
    """
    Assemble slide images into a PPTX file.

    Each image is placed as a full-bleed picture on a blank slide (16:9).
    Returns the PPTX file as bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Emu

    prs = Presentation()

    # Set widescreen 16:9 dimensions (13.333 x 7.5 inches = standard widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout (index 6)
    blank_layout = prs.slide_layouts[6]

    for img_bytes in slide_images:
        slide = prs.slides.add_slide(blank_layout)

        # Add image as full-bleed (covering entire slide)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(
            img_stream,
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# ── Public async API ──────────────────────────────────────────────────

async def generate_presentation(topic: str) -> tuple[bytes, list[bytes]]:
    """
    Generate a 2-slide presentation about the given topic.

    Returns: (pptx_bytes, slide_images) where slide_images is a list of PNG bytes.
    """
    logger.info("Generating presentation: topic=%r", topic)

    # Generate both slide images
    title_prompt = build_title_slide_prompt(topic)
    content_prompt = build_content_slide_prompt(topic)

    # Generate slides sequentially (Gemini rate limits)
    logger.info("Generating title slide...")
    title_image = await asyncio.to_thread(_sync_generate_slide_image, title_prompt)
    if not title_image:
        raise RuntimeError("Failed to generate title slide image from Gemini")

    logger.info("Generating content slide...")
    content_image = await asyncio.to_thread(_sync_generate_slide_image, content_prompt)
    if not content_image:
        raise RuntimeError("Failed to generate content slide image from Gemini")

    slide_images = [title_image, content_image]

    # Assemble into PPTX
    logger.info("Assembling PPTX...")
    pptx_bytes = await asyncio.to_thread(
        _sync_assemble_pptx, slide_images, topic
    )

    logger.info(
        "Presentation ready: %d slides, %.1f KB",
        len(slide_images),
        len(pptx_bytes) / 1024,
    )

    return pptx_bytes, slide_images
